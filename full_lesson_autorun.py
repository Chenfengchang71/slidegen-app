
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from PIL import Image
from io import BytesIO
import openai
import requests

# === CONFIG ===
openai.api_key = os.getenv("OPENAI_API_KEY")  # Set via environment variablegsSS51E-S-G7OWrVUVMBPkWjgwbG4Po-tx1Kztiea-yRvIMYCJEsZ_pENa8vLYP08NUT9n9l7OT3BlbkFJMD-bRXbMED6DWazer-ntPe5ys64FttG0jZo-Z4ZBB57ZasDCOab76WB9S2RXINgy2I2z6LyBkA"
TEMPLATE_PATH = "/mnt/data/generate_ppt.pptx"
OUTPUT_DIR = "/mnt/data/"
FONT_NAME = "Arial"
FONT_SIZE = Pt(24)
SLIDE_WIDTH, SLIDE_HEIGHT = Inches(10), Inches(7.5)
IMAGE_HEIGHT = Inches(3)
IMAGE_WIDTH = Inches(4.25)

# === AI HELPERS ===
def ask_gpt(prompt):
    response = openai.ChatCompletion.create(
        model="gpt-4",
        messages=[{"role": "user", "content": prompt}]
    )
    return response.choices[0].message.content.strip()

# === VOCABULARY WORKFLOW ===
def generate_vocab_words(topic: str) -> list:
    words_text = ask_gpt(f"List 6 important vocabulary words for a lesson on {topic}.")
    words = words_text.strip().split("\n")
    return [w.split(". ")[-1].strip() for w in words][:6]

def generate_image_for_word(word: str) -> Image.Image:
    dalle_response = openai.Image.create(prompt=f"{word}, realistic photo", n=1, size="512x512")
    image_url = dalle_response['data'][0]['url']
    img_data = requests.get(image_url).content
    return Image.open(BytesIO(img_data))

def create_vocab_slides(topic: str, vocab_words: list, image_list: list):
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    for i in range(0, 6, 2):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        title = slide.shapes.title
        title.text = "Vocabulary"
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.size = FONT_SIZE
        title.text_frame.paragraphs[0].font.name = FONT_NAME

        for col, offset in enumerate([Inches(1), Inches(5.25)]):
            word = vocab_words[i + col]
            image = image_list[i + col]
            img_path = f"{OUTPUT_DIR}temp_{i+col}.png"
            image.save(img_path)

            txBox = slide.shapes.add_textbox(offset, Inches(1.2), IMAGE_WIDTH, Inches(0.6))
            tf = txBox.text_frame
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = word
            run.font.bold = True
            run.font.size = FONT_SIZE
            run.font.name = FONT_NAME
            p.alignment = 1

            slide.shapes.add_picture(img_path, offset, Inches(2), width=IMAGE_WIDTH, height=IMAGE_HEIGHT)

    vocab_path = os.path.join(OUTPUT_DIR, f"vocab_slides_{topic.lower().replace(' ', '_')}.pptx")
    prs.save(vocab_path)
    print(f"Vocabulary slides saved to {vocab_path}")

# === LESSON SLIDE WORKFLOW ===
def generate_lesson_content(topic):
    return {
        "title": f"{topic} Lesson",
        "goals": ask_gpt(f"List 3 learning goals for a lesson on {topic}."),
        "warmup": ask_gpt(f"Write 3 warm-up discussion questions for a lesson on {topic}."),
        "reading": ask_gpt(f"Write a short reading passage (5-6 sentences) about {topic}."),
        "followup": ask_gpt(f"Write 3 follow-up questions based on a reading passage about {topic}."),
        "discussion": ask_gpt(f"Write 3 deeper discussion questions about {topic}."),
        "recap": ask_gpt(f"Summarize 3 key recap points from a lesson on {topic}."),
    }

def fill_lesson_template(content, topic):
    prs = Presentation(TEMPLATE_PATH)

    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text = shape.text_frame.text
            if "INSERT_TOPIC_HERE" in text:
                shape.text_frame.text = content["title"]
            elif "INSERT_GOAL_" in text:
                shape.text_frame.text = content["goals"]
            elif "INSERT_QUESTION_" in text:
                shape.text_frame.text = content["warmup"]
            elif "INSERT_READING_TEXT_HERE" in text:
                shape.text_frame.text = content["reading"]
            elif "INSERT_FOLLOWUP_QUESTION_" in text:
                shape.text_frame.text = content["followup"]
            elif "INSERT_DISCUSSION_QUESTION_" in text:
                shape.text_frame.text = content["discussion"]
            elif "INSERT_RECAP_POINT_" in text:
                shape.text_frame.text = content["recap"]

    output_path = os.path.join(OUTPUT_DIR, f"lesson_{topic.lower().replace(' ', '_')}.pptx")
    prs.save(output_path)
    print(f"Lesson slides saved to {output_path}")

# === MASTER RUN ===
def full_lesson_run():
    topic = input("Enter a lesson topic: ")
    # Vocabulary
    vocab_words = generate_vocab_words(topic)
    images = [generate_image_for_word(word) for word in vocab_words]
    create_vocab_slides(topic, vocab_words, images)

    # Non-Vocab Lesson Slides
    lesson_content = generate_lesson_content(topic)
    fill_lesson_template(lesson_content, topic)

if __name__ == "__main__":
    full_lesson_run()
