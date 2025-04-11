import streamlit as st
from pptx import Presentation
from pptx.util import Inches
import io

st.set_page_config(page_title="SlideGen", layout="wide")

st.title("🚀 SlideGen App with Real PowerPoint Export")

# Upload section
uploaded_file = st.file_uploader("📄 Upload a topic file (TXT only for now)", type=["txt"])

# Options
font = st.selectbox("🖋️ Choose a font", ["Arial", "Times New Roman", "Comic Sans"])
theme = st.radio("🎨 Select theme", ["Light", "Dark", "Colorful"])

# Generate slides
if st.button("🎬 Generate Slides"):
    if uploaded_file is not None:
        content = uploaded_file.read().decode("utf-8")
        lines = content.split("\n")

        # Create PowerPoint presentation
        prs = Presentation()
        for i, line in enumerate(lines[:10], 1):
            slide = prs.slides.add_slide(prs.slide_layouts[5])  # blank layout
            title_shape = slide.shapes.title
            textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(5))
            tf = textbox.text_frame
           p = tf.paragraphs[0]
run = p.runs[0]
run.text = line
run.font.name = 'Calibri'


        # Export as .pptx
        pptx_io = io.BytesIO()
        prs.save(pptx_io)
        pptx_io.seek(0)

        st.success("✅ Slides created! Click below to download:")
        st.download_button(
            label="📥 Download Slides",
            data=pptx_io,
            file_name="generated_slides.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )
    else:
        st.error("Please upload a file first.")
